import fitz  # PyMuPDF
import pandas as pd
import docx

from pptx import Presentation

import easyocr
from PIL import Image
import pytesseract

import tempfile
import os


# ==========================================
# PDF
# ==========================================

def read_pdf(file_path):

    pages = []

    pdf = fitz.open(file_path)

    for page_number, page in enumerate(pdf, start=1):

        text = page.get_text()

        print(f"\n===== PAGE {page_number} =====")
        print(text[:500])   

        if text.strip():

            pages.append({
                "page": page_number,
                "text": text
            })

    pdf.close()

    return pages


# ==========================================
# IMAGE (OCR)
# ==========================================

def read_image(file_path):

    text = ""

    try:

        reader = easyocr.Reader(["en", "ar"])

        result = reader.readtext(file_path, detail=0)

        text = "\n".join(result)

    except:

        try:

            text = pytesseract.image_to_string(
                Image.open(file_path),
                lang="eng+ara"
            )

        except:

            text = ""

    return [{
        "page": 1,
        "text": text
    }]


# ==========================================
# DOCX
# ==========================================

def read_docx(file_path):

    doc = docx.Document(file_path)

    text = []

    for p in doc.paragraphs:

        if p.text.strip():

            text.append(p.text)

    for table in doc.tables:

        for row in table.rows:

            cells = [cell.text for cell in row.cells]

            text.append(" | ".join(cells))

    return [{
        "page": 1,
        "text": "\n".join(text)
    }]


# ==========================================
# PPTX
# ==========================================

def read_pptx(file_path):

    prs = Presentation(file_path)

    slides = []

    for slide_number, slide in enumerate(prs.slides, start=1):

        slide_text = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                if shape.text.strip():

                    slide_text.append(shape.text)

        slides.append({

            "page": slide_number,

            "text": "\n".join(slide_text)

        })

    return slides


# ==========================================
# CSV
# ==========================================

def read_csv(file_path):

    df = pd.read_csv(file_path)

    return [{

        "page": 1,

        "text": df.to_string()

    }]


# ==========================================
# MAIN PARSER
# ==========================================

def parse_document(file_path):

    extension = file_path.split(".")[-1].lower()

    if extension == "pdf":

        return read_pdf(file_path)

    elif extension in ["png", "jpg", "jpeg"]:

        return read_image(file_path)

    elif extension == "docx":

        return read_docx(file_path)

    elif extension == "pptx":

        return read_pptx(file_path)

    elif extension == "csv":

        return read_csv(file_path)

    else:

        raise Exception("Unsupported file type.")